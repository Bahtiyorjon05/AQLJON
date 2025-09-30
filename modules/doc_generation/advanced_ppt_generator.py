import asyncio
import tempfile
import logging
import time
import random
import re
import json
from typing import Tuple, Dict, Any, Optional, List
from telegram import Update
from telegram.ext import ContextTypes
from telegram.constants import ParseMode
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from modules.doc_generation.base_generator import BaseDocumentGenerator
from modules.utils import safe_reply, send_typing
from modules.config import Config

logger = logging.getLogger(__name__)

class AdvancedPPTGenerator(BaseDocumentGenerator):
    """Enhanced PowerPoint generator with topic-based styling and advanced visual design."""
    
    def __init__(self, gemini_model, memory_manager):
        super().__init__(gemini_model, memory_manager)
        self.topic_structures = self._setup_topic_structures()
    
    def _setup_topic_structures(self) -> Dict[str, Dict[str, Any]]:
        """Setup topic structures with unique styles, themes, and colors for each category."""
        return {
            'business': {
                'name': 'Business & Corporate',
                'prompt': "Create a professional business presentation about {topic}. Include executive summaries, market analysis, strategic recommendations, financial projections, and implementation timelines. Structure it with clear sections and data-driven insights.",
                'slide_structure': [
                    {'type': 'title', 'title': '{topic}', 'subtitle': 'Professional Business Analysis'},
                    {'type': 'agenda', 'title': 'Agenda'},
                    {'type': 'content', 'title': 'Executive Summary'},
                    {'type': 'content', 'title': 'Market Analysis'},
                    {'type': 'content', 'title': 'Strategic Recommendations'},
                    {'type': 'data', 'title': 'Financial Projections'},
                    {'type': 'timeline', 'title': 'Implementation Timeline'},
                    {'type': 'conclusion', 'title': 'Conclusion & Next Steps'}
                ],
                'visual_elements': {
                    'background': 'gradient',
                    'slide_design': 'corporate',
                    'accent_shapes': ['rectangle', 'chevron'],
                    'color_palette': 'professional',
                    'font_family': 'Calibri',
                    'font_styles': {'bold': True, 'italic': False},
                    'spacing': 'wide',
                    'borders': 'sharp',
                    'corners': 'rounded'
                }
            },
            'technology': {
                'name': 'Technology & Innovation',
                'prompt': "Create a cutting-edge technology presentation about {topic}. Include innovation highlights, technical specifications, implementation roadmap, benefits, challenges, and future outlook. Make it visually engaging with tech-inspired elements.",
                'slide_structure': [
                    {'type': 'title', 'title': '{topic}', 'subtitle': 'Technology Innovation'},
                    {'type': 'agenda', 'title': 'Tech Agenda'},
                    {'type': 'content', 'title': 'Innovation Highlights'},
                    {'type': 'content', 'title': 'Technical Specifications'},
                    {'type': 'roadmap', 'title': 'Implementation Roadmap'},
                    {'type': 'content', 'title': 'Benefits & Challenges'},
                    {'type': 'content', 'title': 'Future Outlook'},
                    {'type': 'conclusion', 'title': 'Tech Conclusion'}
                ],
                'visual_elements': {
                    'background': 'circuit',
                    'slide_design': 'modern',
                    'accent_shapes': ['circle', 'ISOSCELES_TRIANGLE'],
                    'color_palette': 'modern',
                    'font_family': 'Consolas',
                    'font_styles': {'bold': True, 'italic': True},
                    'spacing': 'compact',
                    'borders': 'glow',
                    'corners': 'sharp'
                }
            },
            'education': {
                'name': 'Education & Training',
                'prompt': "Create an educational presentation about {topic}. Include learning objectives, key concepts, practical examples, interactive elements, assessment methods, and resources. Structure it for optimal student engagement and knowledge retention.",
                'slide_structure': [
                    {'type': 'title', 'title': '{topic}', 'subtitle': 'Educational Learning'},
                    {'type': 'agenda', 'title': 'Learning Objectives'},
                    {'type': 'content', 'title': 'Key Concepts'},
                    {'type': 'examples', 'title': 'Practical Examples'},
                    {'type': 'interactive', 'title': 'Interactive Activities'},
                    {'type': 'content', 'title': 'Assessment Methods'},
                    {'type': 'resources', 'title': 'Learning Resources'},
                    {'type': 'conclusion', 'title': 'Summary & Q&A'}
                ],
                'visual_elements': {
                    'background': 'chalkboard',
                    'slide_design': 'educational',
                    'accent_shapes': ['RECTANGLE', 'RECTANGLE'],
                    'color_palette': 'elegant',
                    'font_family': 'Georgia',
                    'font_styles': {'bold': False, 'italic': True},
                    'spacing': 'normal',
                    'borders': 'dashed',
                    'corners': 'rounded'
                }
            },
            'marketing': {
                'name': 'Marketing & Branding',
                'prompt': "Create a dynamic marketing presentation about {topic}. Include campaign overview, target audience analysis, marketing strategies, creative concepts, budget allocation, success metrics, and ROI projections. Make it visually appealing and persuasive.",
                'slide_structure': [
                    {'type': 'title', 'title': '{topic}', 'subtitle': 'Marketing Excellence'},
                    {'type': 'agenda', 'title': 'Campaign Overview'},
                    {'type': 'content', 'title': 'Target Audience Analysis'},
                    {'type': 'content', 'title': 'Marketing Strategies'},
                    {'type': 'content', 'title': 'Creative Concepts'},
                    {'type': 'data', 'title': 'Budget & Metrics'},
                    {'type': 'content', 'title': 'ROI Projections'},
                    {'type': 'conclusion', 'title': 'Next Steps'}
                ],
                'visual_elements': {
                    'background': 'gradient',
                    'slide_design': 'vibrant',
                    'accent_shapes': ['STAR_5_POINT', 'RECTANGLE'],
                    'color_palette': 'vibrant',
                    'font_family': 'Arial',
                    'font_styles': {'bold': True, 'italic': False},
                    'spacing': 'wide',
                    'borders': 'thick',
                    'corners': 'rounded'
                }
            },
            'health': {
                'name': 'Healthcare & Wellness',
                'prompt': "Create a healthcare presentation about {topic}. Include medical overview, health benefits, treatment options, prevention strategies, patient care guidelines, and wellness recommendations. Ensure medical accuracy and compassionate tone.",
                'slide_structure': [
                    {'type': 'title', 'title': '{topic}', 'subtitle': 'Health & Wellness'},
                    {'type': 'agenda', 'title': 'Healthcare Agenda'},
                    {'type': 'content', 'title': 'Medical Overview'},
                    {'type': 'content', 'title': 'Health Benefits'},
                    {'type': 'content', 'title': 'Treatment Options'},
                    {'type': 'content', 'title': 'Prevention Strategies'},
                    {'type': 'content', 'title': 'Patient Care Guidelines'},
                    {'type': 'conclusion', 'title': 'Wellness Recommendations'}
                ],
                'visual_elements': {
                    'background': 'gradient',
                    'slide_design': 'clean',
                    'accent_shapes': ['HEART', 'MATH_PLUS'],
                    'color_palette': 'professional',
                    'font_family': 'Verdana',
                    'font_styles': {'bold': True, 'italic': False},
                    'spacing': 'normal',
                    'borders': 'thin',
                    'corners': 'soft'
                }
            },
            'environment': {
                'name': 'Environment & Sustainability',
                'prompt': "Create an environmental presentation about {topic}. Include ecological impact, sustainability initiatives, green technologies, conservation efforts, policy recommendations, and future sustainability goals. Emphasize eco-friendly practices and environmental responsibility.",
                'slide_structure': [
                    {'type': 'title', 'title': '{topic}', 'subtitle': 'Environmental Sustainability'},
                    {'type': 'agenda', 'title': 'Sustainability Agenda'},
                    {'type': 'content', 'title': 'Ecological Impact'},
                    {'type': 'content', 'title': 'Sustainability Initiatives'},
                    {'type': 'content', 'title': 'Green Technologies'},
                    {'type': 'content', 'title': 'Conservation Efforts'},
                    {'type': 'content', 'title': 'Policy Recommendations'},
                    {'type': 'conclusion', 'title': 'Future Goals'}
                ],
                'visual_elements': {
                    'background': 'leaf_pattern',
                    'slide_design': 'natural',
                    'accent_shapes': ['TEAR', 'RECTANGLE'],
                    'color_palette': 'elegant',
                    'font_family': 'Times New Roman',
                    'font_styles': {'bold': False, 'italic': True},
                    'spacing': 'wide',
                    'borders': 'natural',
                    'corners': 'organic'
                }
            },
            'sports': {
                'name': 'Sports & Recreation',
                'prompt': "Create an energetic sports presentation about {topic}. Include performance analysis, training techniques, team strategies, equipment recommendations, nutrition guidelines, and performance metrics. Make it motivational and action-oriented.",
                'slide_structure': [
                    {'type': 'title', 'title': '{topic}', 'subtitle': 'Sports Excellence'},
                    {'type': 'agenda', 'title': 'Game Plan'},
                    {'type': 'content', 'title': 'Performance Analysis'},
                    {'type': 'content', 'title': 'Training Techniques'},
                    {'type': 'content', 'title': 'Team Strategies'},
                    {'type': 'content', 'title': 'Equipment & Nutrition'},
                    {'type': 'data', 'title': 'Performance Metrics'},
                    {'type': 'conclusion', 'title': 'Championship Mindset'}
                ],
                'visual_elements': {
                    'background': 'ball_pattern',
                    'slide_design': 'energetic',
                    'accent_shapes': ['OVAL', 'RECTANGLE'],
                    'color_palette': 'energizing',
                    'font_family': 'Impact',
                    'font_styles': {'bold': True, 'italic': False},
                    'spacing': 'compact',
                    'borders': 'dynamic',
                    'corners': 'sharp'
                }
            },
            'travel': {
                'name': 'Travel & Tourism',
                'prompt': "Create an exciting travel presentation about {topic}. Include destination highlights, cultural insights, travel itineraries, accommodation options, local experiences, and travel tips. Make it visually stunning and inspiring for travelers.",
                'slide_structure': [
                    {'type': 'title', 'title': '{topic}', 'subtitle': 'Travel Adventure'},
                    {'type': 'agenda', 'title': 'Journey Overview'},
                    {'type': 'content', 'title': 'Destination Highlights'},
                    {'type': 'content', 'title': 'Cultural Insights'},
                    {'type': 'itinerary', 'title': 'Travel Itinerary'},
                    {'type': 'content', 'title': 'Accommodation & Experiences'},
                    {'type': 'content', 'title': 'Local Tips & Tricks'},
                    {'type': 'conclusion', 'title': 'Memorable Experiences'}
                ],
                'visual_elements': {
                    'background': 'airplane_pattern',
                    'slide_design': 'exploratory',
                    'accent_shapes': ['RECTANGLE', 'OVAL'],
                    'color_palette': 'vibrant',
                    'font_family': 'Trebuchet MS',
                    'font_styles': {'bold': True, 'italic': True},
                    'spacing': 'adventurous',
                    'borders': 'exploratory',
                    'corners': 'aerodynamic'
                }
            },
            'finance': {
                'name': 'Finance & Investment',
                'prompt': "Create a comprehensive finance presentation about {topic}. Include market analysis, investment strategies, risk assessment, portfolio recommendations, financial planning, and wealth management insights. Ensure data accuracy and professional presentation.",
                'slide_structure': [
                    {'type': 'title', 'title': '{topic}', 'subtitle': 'Financial Intelligence'},
                    {'type': 'agenda', 'title': 'Financial Agenda'},
                    {'type': 'content', 'title': 'Market Analysis'},
                    {'type': 'content', 'title': 'Investment Strategies'},
                    {'type': 'risk', 'title': 'Risk Assessment'},
                    {'type': 'content', 'title': 'Portfolio Recommendations'},
                    {'type': 'content', 'title': 'Financial Planning'},
                    {'type': 'conclusion', 'title': 'Wealth Management'}
                ],
                'visual_elements': {
                    'background': 'gradient',
                    'slide_design': 'professional',
                    'accent_shapes': ['CHART_X', 'DONUT'],
                    'color_palette': 'professional',
                    'font_family': 'Calibri',
                    'font_styles': {'bold': True, 'italic': False},
                    'spacing': 'balanced',
                    'borders': 'professional',
                    'corners': 'precise'
                }
            },
            'economics': {
                'name': 'Economics & Development',
                'prompt': "Create a comprehensive economics presentation about {topic}. Include economic analysis, development strategies, policy recommendations, market trends, investment opportunities, and growth projections. Make it data-driven with clear economic insights.",
                'slide_structure': [
                    {'type': 'title', 'title': '{topic}', 'subtitle': 'Economic Analysis & Development'},
                    {'type': 'agenda', 'title': 'Economic Overview'},
                    {'type': 'content', 'title': 'Economic Analysis'},
                    {'type': 'content', 'title': 'Development Strategies'},
                    {'type': 'content', 'title': 'Policy Recommendations'},
                    {'type': 'data', 'title': 'Market Trends & Data'},
                    {'type': 'content', 'title': 'Investment Opportunities'},
                    {'type': 'conclusion', 'title': 'Growth Projections'}
                ],
                'visual_elements': {
                    'background': 'gradient',
                    'slide_design': 'professional',
                    'accent_shapes': ['CHART_X', 'RECTANGLE'],
                    'color_palette': 'corporate',
                    'font_family': 'Calibri',
                    'font_styles': {'bold': True, 'italic': False},
                    'spacing': 'balanced',
                    'borders': 'professional',
                    'corners': 'precise'
                }
            },
            'teachers': {
                'name': 'Teachers & Education',
                'prompt': "Create an inspiring presentation about {topic} for teachers. Include educational methodologies, teaching strategies, classroom management techniques, student engagement approaches, professional development opportunities, and educational resources. Make it practical and motivational for educators.",
                'slide_structure': [
                    {'type': 'title', 'title': '{topic}', 'subtitle': 'Educational Excellence for Teachers'},
                    {'type': 'agenda', 'title': 'Teaching Excellence'},
                    {'type': 'content', 'title': 'Educational Methodologies'},
                    {'type': 'content', 'title': 'Teaching Strategies'},
                    {'type': 'content', 'title': 'Classroom Management'},
                    {'type': 'content', 'title': 'Student Engagement'},
                    {'type': 'content', 'title': 'Professional Development'},
                    {'type': 'conclusion', 'title': 'Educational Resources'}
                ],
                'visual_elements': {
                    'background': 'chalkboard',
                    'slide_design': 'educational',
                    'accent_shapes': ['RECTANGLE', 'OVAL'],
                    'color_palette': 'elegant',
                    'font_family': 'Georgia',
                    'font_styles': {'bold': False, 'italic': True},
                    'spacing': 'normal',
                    'borders': 'dashed',
                    'corners': 'rounded'
                }
            },
            'default': {
                'name': 'General Presentation',
                'prompt': "Create a well-structured, engaging presentation about {topic}. Include an introduction, problem statement or context, key insights, detailed analysis, practical examples or case studies, actionable recommendations, and a strong conclusion. Make it visually appealing with relevant emojis and professional formatting. All content should be in the SAME LANGUAGE as the user's input.",
                'slide_structure': [
                    {'type': 'title', 'title': '{topic}', 'subtitle': 'Professional Presentation'},
                    {'type': 'agenda', 'title': 'Presentation Overview'},
                    {'type': 'content', 'title': 'Introduction & Context'},
                    {'type': 'content', 'title': 'Key Insights & Analysis'},
                    {'type': 'examples', 'title': 'Practical Examples'},
                    {'type': 'content', 'title': 'Recommendations'},
                    {'type': 'conclusion', 'title': 'Key Takeaways & Next Steps'}
                ],
                'visual_elements': {
                    'background': 'gradient',
                    'slide_design': 'professional',
                    'accent_shapes': ['RECTANGLE', 'OVAL'],
                    'color_palette': 'awesome',  # Use awesome color scheme with blue, green, emerald, gold
                    'font_family': 'Segoe UI',  # Better default font for presentations
                    'font_styles': {'bold': True, 'italic': False},
                    'spacing': 'balanced',
                    'borders': 'professional',
                    'corners': 'rounded'
                }
            }
        }
    
    def _detect_topic_category(self, topic: str) -> str:
        """Detect the topic category using Gemini AI for intelligent multilingual detection."""
        try:
            # Create a prompt for Gemini to classify the topic
            classification_prompt = f"""
            Analyze the following topic and classify it into ONE of these categories:
            business, technology, education, marketing, health, environment, sports, travel, finance, economics, teachers, default
            
            Topic: "{topic}"
            
            Respond with ONLY the category name in lowercase, nothing else.
            If the topic doesn't clearly fit any category, respond with "default".
            """
            
            # Use a simpler approach for topic detection - let Gemini decide
            # This is more reliable than keyword matching for multilingual content
            response = self.model.generate_content(classification_prompt)
            
            if response and response.text:
                detected_category = response.text.strip().lower()
                # Validate that the detected category is in our known categories
                known_categories = set(self.topic_structures.keys())
                if detected_category in known_categories:
                    return detected_category
            
            # Fallback to default if Gemini response is invalid
            return 'default'
            
        except Exception as e:
            logger.warning(f"Gemini topic detection failed: {e}")
            # Fallback to keyword-based detection as a backup
            return self._detect_topic_category_keyword_fallback(topic)
    
    def _detect_topic_category_keyword_fallback(self, topic: str) -> str:
        """Fallback keyword-based topic detection for when Gemini fails."""
        topic_lower = topic.lower()
        
        # Multilingual keyword mappings for topic detection
        # Note: Order matters - more specific categories should come first
        topic_keywords = {
            'economics': [
                'economics', 'economic', 'development', 'growth', 'policy', 'market',
                'iqtisodiyot', 'iqtisodiy', 'yuksaltirish', 'rivojlanish', 'siyosat', 'bozor',
                'экономика', 'экономический', 'развитие', 'рост', 'политика', 'рынок',
                'milliy bank', 'milliy iqtisodiyot', 'sanoat', 'qishloq xo\'jaligi'
            ],
            'business': [
                'business', 'corporate', 'company', 'enterprise', 'startup', 'entrepreneurship',
                'biznes', 'kompaniya', 'korxona', 'tadbirkorlik', 'boshlash',
                'бизнес', 'компания', 'предприятие', 'стартап', 'предпринимательство'
            ],
            'technology': [
                'technology', 'tech', 'innovation', 'digital', 'software', 'ai', 'artificial intelligence',
                'texnologiya', 'innovatsiya', 'raqamli', 'dasturiy', 'sun\'iy intellekt',
                'технология', 'инновация', 'цифровой', 'программное обеспечение', 'искусственный интеллект'
            ],
            'education': [
                'education', 'learning', 'teaching', 'training', 'school', 'university', 'course',
                'ta\'lim', 'o\'qish', 'o\'qitish', 'kurs', 'maktab', 'universitet',
                'образование', 'обучение', 'преподавание', 'тренинг', 'школа', 'университет'
            ],
            'marketing': [
                'marketing', 'advertising', 'branding', 'campaign', 'promotion', 'social media',
                'marketing', 'reklama', 'brend', 'kampaniya', 'ijtimoiy tarmoq',
                'маркетинг', 'реклама', 'брендинг', 'кампания', 'продвижение', 'социальные сети'
            ],
            'health': [
                'health', 'medical', 'wellness', 'fitness', 'nutrition', 'medicine', 'therapy',
                'salomatlik', 'tibbiy', 'fitnes', 'ovqatlanish', 'dorilar', 'terapiya',
                'здоровье', 'медицинский', 'благополучие', 'фитнес', 'питание', 'медицина'
            ],
            'environment': [
                'environment', 'sustainability', 'green', 'ecology', 'climate', 'conservation',
                'atrof muhit', 'barqarorlik', 'yashil', 'ekologiya', 'iqlim', 'himoya',
                'окружающая среда', 'устойчивость', 'зеленый', 'экология', 'климат', 'охрана',
                'atmosfer', 'havo', 'tabiat'
            ],
            'sports': [
                'sports', 'athletics', 'fitness', 'competition', 'team', 'game', 'exercise',
                'sport', 'atletika', 'fitnes', 'musobaqa', 'jamoa', 'o\'yin',
                'спорт', 'атлетика', 'фитнес', 'соревнование', 'команда', 'игра',
                'sportning ahamiyati', 'sport haqida', 'sportning foydasi'
            ],
            'travel': [
                'travel', 'tourism', 'destination', 'vacation', 'trip', 'journey', 'adventure',
                'sayohat', 'turizm', 'manzil', 'dam olish', 'safar', 'sarguzasht',
                'путешествие', 'туризм', 'назначение', 'отпуск', 'поездка', 'приключение'
            ],
            'finance': [
                'finance', 'investment', 'money', 'banking', 'economy', 'wealth', 'budget',
                'moliya', 'investitsiya', 'pul', 'bank', 'iqtisodiyot', 'boylik', 'byudjet',
                'қарз', 'қарз олиш', 'қарз берish'
            ],
            'teachers': [
                'teachers', 'teacher', 'educator', 'education', 'teaching', 'instructor',
                'oqituvchi', 'oqituvchilar', 'pedagog', 'ta\'lim', 'o\'qitish',
                'учитель', 'учителя', 'педагог', 'образование', 'преподавание',
                'ustoz', 'ustozlar', 'maktab oqituvchisi'
            ]
        }
        
        # Check for matches in each category
        for category, keywords in topic_keywords.items():
            for keyword in keywords:
                if keyword in topic_lower:
                    return category
        
        # Default category if no match found
        return 'default'
    
    async def generate(self, update: Update, context: ContextTypes.DEFAULT_TYPE, topic: str, content_context: str = ""):
        """Generate an enhanced PowerPoint presentation with topic-specific styling."""
        processing_msg = None
        temp_file = None
        
        try:
            # Send initial processing message with immediate response
            try:
                from modules.utils import send_fast_reply
                if update.message:
                    send_fast_reply(update.message, "<b>📊 PowerPoint taqdimotni tuzyapman. Ozgina kutib turing... ⏳</b>")
                    # Send typing indicator
                    await send_typing(update)
            except:
                # Fallback if fast reply fails
                processing_msg = await self._send_processing_message(update, f"<b>📊 PowerPoint taqdimotni tuzyapman. Ozgina kutib turing... ⏳</b>")
                # Send typing indicator
                await send_typing(update)
            
            # Track document generation
            await self._track_document_generation(update, "powerpoint")
            
            # Detect topic category for styling with timeout to prevent blocking
            try:
                topic_category = await asyncio.wait_for(
                    asyncio.to_thread(self._detect_topic_category, topic),
                    timeout=15.0  # Increased timeout to 15 seconds
                )
            except asyncio.TimeoutError:
                logger.warning("Topic detection timed out, using default category")
                topic_category = 'default'
            except Exception as e:
                logger.warning(f"Topic detection failed: {e}, using default category")
                topic_category = 'default'
            
            structure = self.topic_structures.get(topic_category, self.topic_structures['default'])
            
            # Generate content with Gemini with increased timeout and retry logic
            content = None
            for attempt in range(2):  # Try up to 2 times
                try:
                    content = await asyncio.wait_for(
                        self._generate_presentation_content(topic, structure, content_context),
                        timeout=Config.PROCESSING_TIMEOUT * 2  # Use configured timeout
                    )
                    break  # Success, exit retry loop
                except asyncio.TimeoutError:
                    logger.warning(f"Content generation timed out on attempt {attempt + 1}")
                    if attempt == 1:  # Last attempt
                        content = self._create_fallback_content(topic, structure)
                except Exception as e:
                    logger.error(f"Content generation failed on attempt {attempt + 1}: {e}")
                    if attempt == 1:  # Last attempt
                        content = self._create_fallback_content(topic, structure)
            
            # Ensure content is not None
            if content is None:
                content = self._create_fallback_content(topic, structure)
            
            # Create presentation in a separate thread to avoid blocking with timeout
            temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
            temp_file.close()
            
            # Offload presentation creation to a separate thread with timeout
            try:
                await asyncio.wait_for(
                    asyncio.get_event_loop().run_in_executor(
                        None, 
                        self._create_enhanced_presentation, 
                        topic, content, structure, temp_file.name  # content is guaranteed to be a dict now
                    ),
                    timeout=Config.PROCESSING_TIMEOUT  # Use configured timeout
                )
            except asyncio.TimeoutError:
                logger.error("Presentation creation timed out")
                raise Exception("Presentation creation took too long")
            except Exception as e:
                logger.error(f"Presentation creation failed: {e}")
                raise Exception(f"Failed to create presentation: {e}")
            
            # Generate filename
            filename = await self._generate_filename(topic, "powerpoint")
            filename_with_ext = f"{filename}.pptx"
            
            # Send file to user
            with open(temp_file.name, 'rb') as f:
                if update.message:
                    await update.message.reply_document(
                        document=f,
                        filename=filename_with_ext,
                        caption=f"📊 <b>'{filename}' mavzusida professional PowerPoint taqdimot</b>\n📁 Fayl nomi: {filename_with_ext}",
                        parse_mode=ParseMode.HTML
                    )
                elif update.effective_message:
                    # Fallback if message is not available
                    await update.effective_message.reply_document(
                        document=f,
                        filename=filename_with_ext,
                        caption=f"📊 <b>'{filename}' mavzusida professional PowerPoint taqdimot</b>\n📁 Fayl nomi: {filename_with_ext}",
                        parse_mode=ParseMode.HTML
                    )
                else:
                    # Last resort fallback
                    await safe_reply(update, f"📊 <b>'{filename}' mavzusida professional PowerPoint taqdimot</b>\n📁 Fayl nomi: {filename_with_ext}", parse_mode=ParseMode.HTML)
            
            # Send success message
            try:
                from modules.utils import send_fast_reply
                if update.message:
                    send_fast_reply(update.message, f"✅ <b>'{filename}' nomli PowerPoint taqdimotingiz muvaffaqiyatli tuzildi va yuborildi!</b>\n📥 Ajoyib dizayn va did bilan tuzilgan faylingizdan zavqlaning!", parse_mode=ParseMode.HTML)
            except:
                # Fallback if fast reply fails
                await self._send_success_message(
                    processing_msg,
                    f"✅ <b>'{filename}' nomli PowerPoint taqdimotingiz muvaffaqiyatli tuzildi va yuborildi!</b>\n"
                    f"📥 Ajoyib dizayn va did bilan tuzilgan faylingizdan zavqlaning!"
                )
            
        except Exception as e:
            logger.error(f"PowerPoint generation failed: {e}", exc_info=True)
            # Send error message to user
            try:
                from modules.utils import send_fast_reply
                if update.message:
                    send_fast_reply(update.message, "❌ PowerPoint taqdimotini yaratishda xatolik yuz berdi. Iltimos, keyinroq qayta urinib ko'ring.", parse_mode=ParseMode.HTML)
            except:
                # Fallback if fast reply fails
                await self._handle_document_generation_error(
                    processing_msg,
                    "❌ PowerPoint taqdimotini yaratishda xatolik yuz berdi. Iltimos, keyinroq qayta urinib ko'ring."
                )
        finally:
            # Cleanup temporary file
            if temp_file and temp_file.name:
                try:
                    import os
                    os.unlink(temp_file.name)
                except Exception as e:
                    logger.warning(f"Failed to cleanup temp file: {e}")
    
    async def _generate_presentation_content(self, topic: str, structure: Dict[str, Any], content_context: str = "") -> Dict[str, Any]:
        """Generate presentation content using Gemini AI."""
        try:
            # Create detailed prompt for Gemini with clearer instructions
            prompt = structure['prompt'].format(topic=topic)
            if content_context:
                prompt += f"\n\nAdditional context: {content_context}"
            
            prompt += """
            
            Please structure your response in the following JSON format:
            {
                "title_slide": {
                    "title": "Main Title",
                    "subtitle": "Subtitle"
                },
                "agenda": ["Point 1", "Point 2", "Point 3", ...],
                "slides": [
                    {
                        "title": "Slide Title",
                        "content": ["Bullet point 1", "Bullet point 2", ...],
                        "type": "content|data|timeline|etc."
                    },
                    ...
                ],
                "conclusion": {
                    "title": "Conclusion Title",
                    "content": ["Summary point 1", "Summary point 2", ...]
                }
            }
            
            Requirements:
            1. All content should be in the SAME LANGUAGE as the user's input
            2. Include relevant emojis where appropriate
            3. Keep bullet points concise but informative
            4. Ensure content matches the slide structure defined
            5. Make it engaging and professionally formatted
            6. Create at least 5 content slides for comprehensive coverage
            7. Each slide should have 3-5 bullet points
            8. Use clear, professional language
            """
            
            # Generate content with Gemini with increased timeout
            try:
                response = await asyncio.wait_for(
                    asyncio.to_thread(lambda: self.model.generate_content(prompt)),
                    timeout=45.0  # Increased timeout to 45 seconds
                )
            except asyncio.TimeoutError:
                logger.warning("Content generation timed out, retrying with simpler prompt")
                # Try with a simpler prompt if the first attempt times out
                simple_prompt = f"""
                Create a professional presentation about '{topic}' in the SAME LANGUAGE as the topic.
                
                Structure it as JSON with these sections:
                - title_slide: with title and subtitle
                - agenda: list of 5-7 main points
                - slides: 5-7 content slides, each with title and 3-5 bullet points
                - conclusion: with title and 3-5 summary points
                
                All content must be in the SAME LANGUAGE as '{topic}'.
                Include relevant emojis.
                Keep it professional and engaging.
                """
                
                response = await asyncio.wait_for(
                    asyncio.to_thread(lambda: self.model.generate_content(simple_prompt)),
                    timeout=Config.PROCESSING_TIMEOUT  # Use configured timeout
                )
            
            if not response or not response.text:
                raise Exception("Empty response from Gemini")
            
            # Extract JSON from response
            import re
            
            # Try to extract JSON from code blocks
            json_match = re.search(r'\{.*\}', response.text, re.DOTALL)
            if json_match:
                json_str = json_match.group(0)
                try:
                    content = json.loads(json_str)
                    # Validate that we have the required structure
                    if isinstance(content, dict) and 'title_slide' in content and 'agenda' in content:
                        return content
                except json.JSONDecodeError:
                    pass
            
            # If JSON extraction fails, create fallback content
            return self._create_fallback_content(topic, structure)
            
        except Exception as e:
            logger.error(f"Content generation failed: {e}")
            return self._create_fallback_content(topic, structure)
    
    def _create_fallback_content(self, topic: str, structure: Dict[str, Any]) -> Dict[str, Any]:
        """Create fallback content when AI generation fails."""
        return {
            "title_slide": {
                "title": topic,
                "subtitle": "Professional Presentation"
            },
            "agenda": [
                "Introduction & Context", 
                "Key Insights & Analysis", 
                "Practical Examples", 
                "Recommendations", 
                "Key Takeaways & Next Steps"
            ],
            "slides": [
                {
                    "title": "Introduction & Context",
                    "content": [
                        f"🔹 Welcome to our comprehensive presentation on {topic}",
                        f"🔹 This presentation provides valuable insights about {topic}",
                        f"🔹 We'll explore key aspects and practical applications of {topic}",
                        f"🔹 Professional analysis and recommendations will be provided"
                    ],
                    "type": "content"
                },
                {
                    "title": "Key Insights & Analysis",
                    "content": [
                        "📊 Important findings and observations",
                        "📈 Data-driven insights and trends",
                        "🔍 Critical factors and considerations",
                        "💡 Key discoveries and implications"
                    ],
                    "type": "content"
                },
                {
                    "title": "Practical Examples",
                    "content": [
                        "💼 Real-world applications and case studies",
                        "🎯 Best practices and proven approaches",
                        "📚 Illustrative examples and scenarios",
                        "🔧 Practical implementation strategies"
                    ],
                    "type": "examples"
                },
                {
                    "title": "Recommendations",
                    "content": [
                        "✅ Actionable strategies and solutions",
                        "📋 Implementation guidelines and steps",
                        "🚀 Next steps and future considerations",
                        "🏆 Best practices for success"
                    ],
                    "type": "content"
                }
            ],
            "conclusion": {
                "title": "Key Takeaways & Next Steps",
                "content": [
                    "⭐ Summary of key insights and findings",
                    "⭐ Recommended actions and approaches",
                    "⭐ Thank you for your attention and engagement",
                    "⭐ Questions and discussion welcome"
                ]
            }
        }
    
    def _create_enhanced_presentation(self, topic: str, content: Dict[str, Any], structure: Dict[str, Any], file_path: str):
        """Create an enhanced presentation with topic-specific styling and save to file."""
        prs = PPTXPresentation()
        
        # Set presentation size
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Get visual elements for this topic
        visual_elements = structure.get('visual_elements', self.topic_structures['default']['visual_elements'])
        color_scheme = self._select_color_scheme_by_topic(topic, "powerpoint")
        
        # Create title slide
        title_slide = content.get('title_slide', {})
        self._create_title_slide(prs, title_slide, visual_elements, color_scheme)
        
        # Create agenda slide
        agenda_items = content.get('agenda', [])
        self._create_agenda_slide(prs, agenda_items, visual_elements, color_scheme)
        
        # Create content slides
        slides = content.get('slides', [])
        for slide_data in slides:
            self._create_content_slide(prs, slide_data, visual_elements, color_scheme)
        
        # Create conclusion slide
        conclusion = content.get('conclusion', {})
        self._create_conclusion_slide(prs, conclusion, visual_elements, color_scheme)
        
        # Save presentation to file
        prs.save(file_path)
    
    def _create_title_slide(self, prs, title_slide: Dict[str, str], 
                           visual_elements: Dict[str, Any], color_scheme: Dict[str, str]):
        """Create a beautifully styled title slide with centered content."""
        slide_layout = prs.slide_layouts[0]  # Title Slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add background design
        self._add_background_design(slide, visual_elements, color_scheme)
        
        # Get visual styling elements
        font_family = visual_elements.get('font_family', 'Segoe UI')  # Better default font
        font_styles = visual_elements.get('font_styles', {'bold': True, 'italic': False})
        
        # Style title with larger font size and center alignment
        title = slide.shapes.title
        if title and hasattr(title, 'text_frame'):
            title.text = title_slide.get('title', 'Presentation')
            title.text_frame.paragraphs[0].font.size = Pt(54)  # Even larger font
            title.text_frame.paragraphs[0].font.name = font_family
            title.text_frame.paragraphs[0].font.bold = font_styles.get('bold', True)
            title.text_frame.paragraphs[0].font.italic = font_styles.get('italic', False)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(
                color_scheme.get('primary', '#2C3E50').lstrip('#')
            )
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Center vertically as well
            title.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # Middle alignment
            
            # Ensure text fits well
            title.width = Inches(10)
            title.height = Inches(2)
            title.left = Inches(1.665)  # Center horizontally
            title.top = Inches(2)  # Position lower for better visual balance
        
        # Style subtitle with larger font size and center alignment
        subtitle_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if subtitle_placeholder and hasattr(subtitle_placeholder, 'text_frame'):
            subtitle_placeholder.text = title_slide.get('subtitle', 'Professional Presentation')
            subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle_placeholder.text_frame.paragraphs[0].font.name = font_family
            subtitle_placeholder.text_frame.paragraphs[0].font.bold = font_styles.get('bold', False)
            subtitle_placeholder.text_frame.paragraphs[0].font.italic = font_styles.get('italic', True)
            subtitle_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(
                color_scheme.get('secondary', '#3498DB').lstrip('#')
            )
            subtitle_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Center vertically as well
            subtitle_placeholder.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # Middle alignment
            
            # Ensure text fits well
            subtitle_placeholder.width = Inches(10)
            subtitle_placeholder.height = Inches(1)
            subtitle_placeholder.left = Inches(1.665)  # Center horizontally
            subtitle_placeholder.top = Inches(4.5)  # Position appropriately

    def _create_agenda_slide(self, prs, agenda_items: List[str],
                            visual_elements: Dict[str, Any], color_scheme: Dict[str, str]):
        """Create an agenda slide with visual enhancements and centered content."""
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add background design
        self._add_background_design(slide, visual_elements, color_scheme)
        
        # Get visual styling elements
        font_family = visual_elements.get('font_family', 'Segoe UI')  # Better default font
        font_styles = visual_elements.get('font_styles', {'bold': True, 'italic': False})
        
        # Style title with larger font size and center alignment
        title = slide.shapes.title
        if title and hasattr(title, 'text_frame'):
            title.text = "Presentation Overview"  # More engaging title
            title.text_frame.paragraphs[0].font.size = Pt(32)
            title.text_frame.paragraphs[0].font.name = font_family
            title.text_frame.paragraphs[0].font.bold = font_styles.get('bold', True)
            title.text_frame.paragraphs[0].font.italic = font_styles.get('italic', False)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(
                color_scheme.get('primary', '#2C3E50').lstrip('#')
            )
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add agenda content with better formatting and centering
        content_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if content_placeholder and hasattr(content_placeholder, 'text_frame'):
            text_frame = content_placeholder.text_frame
            text_frame.clear()  # Clear existing content
            
            # Center the text frame vertically
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # Middle alignment
            
            for i, item in enumerate(agenda_items):
                p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                p.text = f"{i+1}. 📋 {item}"
                p.font.size = Pt(20)
                p.font.name = font_family
                p.font.bold = font_styles.get('bold', True)  # Make agenda items bold
                p.font.italic = font_styles.get('italic', False)
                p.font.color.rgb = RGBColor.from_string(
                    color_scheme.get('text', '#2C3E50').lstrip('#')
                )
                p.space_before = Pt(15)  # Add spacing before each item
                p.space_after = Pt(15)   # Add spacing after each item
                p.alignment = PP_ALIGN.CENTER  # Center align text
                p.level = 0  # No indentation

    def _create_content_slide(self, prs, slide_data: Dict[str, Any],
                             visual_elements: Dict[str, Any], color_scheme: Dict[str, str]):
        """Create a content slide with visual enhancements and centered content."""
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add background design
        self._add_background_design(slide, visual_elements, color_scheme)
        
        # Get visual styling elements
        font_family = visual_elements.get('font_family', 'Segoe UI')  # Better default font
        font_styles = visual_elements.get('font_styles', {'bold': True, 'italic': False})
        
        # Style title with larger font size and center alignment
        title = slide.shapes.title
        if title and hasattr(title, 'text_frame'):
            title.text = slide_data.get('title', 'Slide Title')
            title.text_frame.paragraphs[0].font.size = Pt(28)
            title.text_frame.paragraphs[0].font.name = font_family
            title.text_frame.paragraphs[0].font.bold = font_styles.get('bold', True)
            title.text_frame.paragraphs[0].font.italic = font_styles.get('italic', False)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(
                color_scheme.get('primary', '#2C3E50').lstrip('#')
            )
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add content with better formatting and centering
        content_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if content_placeholder and hasattr(content_placeholder, 'text_frame'):
            text_frame = content_placeholder.text_frame
            text_frame.clear()  # Clear existing content
            
            # Center the text frame vertically
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # Middle alignment
            
            content_items = slide_data.get('content', [])
            slide_type = slide_data.get('type', 'content')
            
            # Add contextual emojis based on slide type
            emoji_map = {
                'data': '📊',
                'timeline': '📅',
                'roadmap': '🛣️',
                'risk': '⚠️',
                'conclusion': '✅',
                'summary': '📌',
                'examples': '💡',
                'interactive': '🎯',
                'resources': '📚',
                'itinerary': '🗺️',
                'content': '🔹'  # Default emoji for content slides
            }
            
            emoji_prefix = emoji_map.get(slide_type, '🔹')
            
            for i, item in enumerate(content_items):
                p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                # Add emoji prefix if not already present
                if not any(item.startswith(emoji) for emoji in ['📊', '📅', '🛣️', '⚠️', '✅', '📌', '💡', '🎯', '📚', '🗺️', '🔹']):
                    p.text = f"{emoji_prefix} {item}"
                else:
                    p.text = item
                p.font.size = Pt(18)
                p.font.name = font_family
                p.font.bold = True  # Make conclusion items bold
                p.font.italic = False
                p.font.color.rgb = RGBColor.from_string(
                    color_scheme.get('text', '#2C3E50').lstrip('#')
                )
                p.space_before = Pt(12)  # Add spacing before each item
                p.space_after = Pt(12)   # Add spacing after each item
                p.alignment = PP_ALIGN.CENTER  # Center align text
                p.level = 0  # No indentation

    def _create_conclusion_slide(self, prs, conclusion: Dict[str, Any],
                                visual_elements: Dict[str, Any], color_scheme: Dict[str, str]):
        """Create a conclusion slide with visual enhancements and centered content."""
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add background design
        self._add_background_design(slide, visual_elements, color_scheme)
        
        # Get visual styling elements
        font_family = visual_elements.get('font_family', 'Segoe UI')  # Better default font
        font_styles = visual_elements.get('font_styles', {'bold': True, 'italic': False})
        
        # Style title with larger font size and center alignment
        title = slide.shapes.title
        if title and hasattr(title, 'text_frame'):
            title.text = conclusion.get('title', 'Key Takeaways & Next Steps')
            title.text_frame.paragraphs[0].font.size = Pt(28)
            title.text_frame.paragraphs[0].font.name = font_family
            title.text_frame.paragraphs[0].font.bold = font_styles.get('bold', True)
            title.text_frame.paragraphs[0].font.italic = font_styles.get('italic', False)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor.from_string(
                color_scheme.get('primary', '#2C3E50').lstrip('#')
            )
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Add content with better formatting and centering
        content_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if content_placeholder and hasattr(content_placeholder, 'text_frame'):
            text_frame = content_placeholder.text_frame
            text_frame.clear()  # Clear existing content
            
            # Center the text frame vertically
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # Middle alignment
            
            content_items = conclusion.get('content', [])
            
            # Add star emoji for key takeaways
            for i, item in enumerate(content_items):
                p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                # Add star emoji if not already present
                if not item.startswith('⭐'):
                    p.text = f"⭐ {item}"
                else:
                    p.text = item
                p.font.size = Pt(18)
                p.font.name = font_family
                p.font.bold = True  # Make conclusion items bold
                p.font.italic = False
                p.font.color.rgb = RGBColor.from_string(
                    color_scheme.get('text', '#2C3E50').lstrip('#')
                )
                p.space_before = Pt(12)  # Add spacing before each item
                p.space_after = Pt(12)   # Add spacing after each item
                p.alignment = PP_ALIGN.CENTER  # Center align text
                p.level = 0  # No indentation
    
    def _add_background_design(self, slide, visual_elements: Dict[str, Any], color_scheme: Dict[str, str]):
        """Add background design elements based on topic category."""
        background_type = visual_elements.get('background', 'gradient')
        
        # Add decorative shapes based on topic - only in corners to keep center clean
        accent_shapes = visual_elements.get('accent_shapes', ['rectangle'])
        for shape_type in accent_shapes:
            self._add_accent_shape(slide, shape_type, color_scheme, visual_elements)
    
    def _add_accent_shape(self, slide, shape_type: str, color_scheme: Dict[str, str], visual_elements: Optional[Dict[str, Any]] = None):
        """Add decorative accent shapes to slides - only in corners to keep center clean."""
        # Position shapes only in corners to avoid interfering with content
        corners = visual_elements.get('corners', 'rounded') if visual_elements else 'rounded'
        
        if corners == 'sharp':
            positions = [
                (Inches(0.3), Inches(0.3)),  # Top-left
                (Inches(12.53), Inches(0.3)),  # Top-right
                (Inches(0.3), Inches(6.7)),  # Bottom-left
                (Inches(12.53), Inches(6.7))  # Bottom-right
            ]
        elif corners == 'organic':
            positions = [
                (Inches(0.7), Inches(0.7)),  # Top-left
                (Inches(12.13), Inches(0.7)),  # Top-right
                (Inches(0.7), Inches(6.3)),  # Bottom-left
                (Inches(12.13), Inches(6.3))  # Bottom-right
            ]
        elif corners == 'aerodynamic':
            positions = [
                (Inches(0.4), Inches(0.6)),  # Top-left
                (Inches(12.43), Inches(0.4)),  # Top-right
                (Inches(0.6), Inches(6.5)),  # Bottom-left
                (Inches(12.23), Inches(6.7))  # Bottom-right
            ]
        else:  # Default rounded corners
            positions = [
                (Inches(0.5), Inches(0.5)),  # Top-left
                (Inches(12.33), Inches(0.5)),  # Top-right
                (Inches(0.5), Inches(6.5)),  # Bottom-left
                (Inches(12.33), Inches(6.5))  # Bottom-right
            ]
        
        shape_mapping = {
            'rectangle': MSO_SHAPE.RECTANGLE,
            'circle': MSO_SHAPE.OVAL,
            'ISOSCELES_TRIANGLE': MSO_SHAPE.ISOSCELES_TRIANGLE,
            'chevron': MSO_SHAPE.CHEVRON,
            'STAR_5_POINT': MSO_SHAPE.STAR_5_POINT,
            'HEART': MSO_SHAPE.HEART,
            'MATH_PLUS': MSO_SHAPE.MATH_PLUS,
            'RECTANGLE': MSO_SHAPE.RECTANGLE,  # Using rectangle as substitute
            'OVAL': MSO_SHAPE.OVAL,  # Using oval as substitute
            'CHART_X': MSO_SHAPE.CHART_X,
            'DONUT': MSO_SHAPE.DONUT,
            'TEAR': MSO_SHAPE.TEAR
        }
        
        shape_enum = shape_mapping.get(shape_type, MSO_SHAPE.RECTANGLE)
        
        # Adjust size based on spacing - keep them small to avoid clutter
        spacing = visual_elements.get('spacing', 'normal') if visual_elements else 'normal'
        if spacing == 'compact':
            size = Inches(0.3)  # Smaller shapes
        elif spacing == 'wide':
            size = Inches(0.5)  # Medium shapes
        else:
            size = Inches(0.4)  # Default size
        
        for x, y in positions:
            try:
                shape = slide.shapes.add_shape(
                    shape_enum,
                    x, y,
                    size, size
                )
                # Style the shape based on border style
                borders = visual_elements.get('borders', 'thin') if visual_elements else 'thin'
                
                # Fill style with transparency to keep content readable
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor.from_string(
                    color_scheme.get('accent', '#F1C40F').lstrip('#')
                )
                
                # Make shape more transparent to avoid interfering with content
                fill.alpha = 0.2  # More transparent
                
                # Border style
                line = shape.line
                if borders == 'thick':
                    line.width = Pt(2)  # Thicker but not too thick
                elif borders == 'glow':
                    line.width = Pt(1.5)
                    # Note: Glow effect not directly supported, using thicker line
                elif borders == 'dashed':
                    line.dash_style = 2  # Dashed line
                    line.width = Pt(1)
                else:  # thin or default
                    line.width = Pt(0.5)  # Thinner lines
                
                line.color.rgb = RGBColor.from_string(
                    color_scheme.get('primary', '#2C3E50').lstrip('#')
                )
                
            except Exception:
                # If shape creation fails, skip it
                pass



    def _get_fallback_content(self, topic: str) -> Dict[str, Any]:
        """Generate fallback content when primary generation fails"""
        return {
            "title_slide": {
                "title": topic,
                "subtitle": "Professional Presentation"
            },
            "agenda": [
                "Introduction",
                "Key Points",
                "Conclusion"
            ],
            "slides": [
                {
                    "title": "Introduction",
                    "content": [
                        f"Welcome to our presentation on {topic}",
                        "This presentation provides an overview of the topic",
                        "We'll cover key aspects and important insights"
                    ],
                    "type": "content"
                },
                {
                    "title": "Key Points",
                    "content": [
                        "Important findings and observations",
                        "Data-driven insights and trends",
                        "Critical factors and considerations"
                    ],
                    "type": "content"
                }
            ],
            "conclusion": {
                "title": "Conclusion",
                "content": [
                    "Summary of key insights and findings",
                    "Recommended actions and approaches",
                    "Thank you for your attention and engagement",
                    "Questions and discussion welcome"
                ]
            }
        }

# Export the class
__all__ = ['AdvancedPPTGenerator']