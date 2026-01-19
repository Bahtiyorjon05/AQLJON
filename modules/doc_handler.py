import asyncio
import tempfile
import time
import logging
import os
import mimetypes
from datetime import datetime
from telegram import Update, Document
from telegram.ext import ContextTypes
from telegram.constants import ParseMode
from modules.utils import safe_reply, send_typing, safe_edit_message
from modules.config import Config
from modules.memory import MemoryManager
from modules.retry_utils import generate_content_with_retry, upload_file_with_retry, wait_for_file_active

# Import libraries for Office document text extraction
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from openpyxl import load_workbook
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)

class DocumentHandler:
    """Handles document processing for the AQLJON bot - supports all file types"""
    
    def __init__(self, gemini_model, memory_manager: MemoryManager):
        self.model = gemini_model
        self.memory = memory_manager
        self.rag_chain = None  # RAG system for semantic search (set via set_rag())
        # Track active document processing tasks per user
        self.active_tasks = {}
        # Cleanup task will be started by main.py after event loop starts
        self._cleanup_task = None

    def set_rag(self, rag_chain):
        """Set RAG chain for document memory integration"""
        self.rag_chain = rag_chain
        logger.info("✅ RAG chain integrated into DocumentHandler")
    
    def _get_user_task_key(self, chat_id, task_id=None):
        """Generate a unique key for tracking user document processing tasks"""
        if task_id:
            return f"{chat_id}_document_{task_id}"
        return f"{chat_id}_document"
    
    def _is_task_active(self, chat_id, task_id=None):
        """Check if a specific document processing task is already active for this user"""
        task_key = self._get_user_task_key(chat_id, task_id)
        return task_key in self.active_tasks and not self.active_tasks[task_key].done()
    
    def _register_task(self, chat_id, task, task_id=None):
        """Register a document processing task for this user"""
        task_key = self._get_user_task_key(chat_id, task_id)
        self.active_tasks[task_key] = task
    
    def _unregister_task(self, chat_id, task_id=None):
        """Unregister a completed document processing task"""
        task_key = self._get_user_task_key(chat_id, task_id)
        if task_key in self.active_tasks:
            del self.active_tasks[task_key]

    async def _cleanup_completed_tasks(self):
        """Periodically clean up completed tasks to prevent memory leaks"""
        while True:
            try:
                await asyncio.sleep(3600)  # Run every hour

                # Find all completed tasks
                completed_keys = [
                    key for key, task in self.active_tasks.items()
                    if task.done()
                ]

                # Remove completed tasks
                for key in completed_keys:
                    del self.active_tasks[key]

                if completed_keys:
                    logger.info(f"🧹 Cleaned up {len(completed_keys)} completed document tasks. Active tasks: {len(self.active_tasks)}")

            except Exception as e:
                logger.error(f"Error in document task cleanup: {e}")

    async def handle_document(self, update: Update, context: ContextTypes.DEFAULT_TYPE):
        """Handle document uploads and analysis - allow concurrent processing"""
        if not update or not update.message or not update.message.document:
            return
            
        chat_id = str(update.effective_chat.id) if update and update.effective_chat else "unknown"
        
        # Generate a unique task ID for this document processing request
        import time
        task_id = str(int(time.time() * 1000))  # Unique timestamp-based ID
        
        await send_typing(update)
        document: Document = update.message.document
        
        # Immediate response - don't block other users
        analyzing_msg = await safe_reply(update,
            "📄 <b>Hujjatni qabul qildim!</b>\n\n"
            "⏳ <i>Tahlil qilyapman... Boshqa buyruqlar yuboravering, men parallel ishlayman!</i>\n\n"
            "📱 <i>Hujjat tahlili tayyor bo'lganda aytaman 😊</i>",
            parse_mode=ParseMode.HTML
        )
        
        # Process document in background with task tracking
        task = asyncio.create_task(self._process_document_background(
            document, chat_id, analyzing_msg, update, context, task_id
        ))
        self._register_task(chat_id, task, task_id)
        
        # Don't await the task - let it run in background
        # Immediately track activity and return - don't block!
        self.memory.track_user_activity(chat_id, "documents", update)
    
    async def _process_document_background(self, document: Document, chat_id: str, analyzing_msg, update: Update, context: ContextTypes.DEFAULT_TYPE, task_id: str):
        """Process document in background without blocking other users"""
        try:
            # Check file size
            if document.file_size and document.file_size > Config.MAX_FILE_SIZE:
                await analyzing_msg.edit_text(
                    "❌ Fayl juda katta. Maksimal hajm: 20MB",
                    parse_mode=ParseMode.HTML
                )
                return
            
            file = await context.bot.get_file(document.file_id)
            
            # Create temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=f"_{document.file_name}") as tmp_file:
                tmp_path = tmp_file.name
            
            try:
                # Download file with timeout
                await asyncio.wait_for(
                    file.download_to_drive(custom_path=tmp_path),
                    timeout=Config.DOWNLOAD_TIMEOUT
                )
                
                # Wait a moment for file to be fully written
                await asyncio.sleep(0.5)
                
                # Check if file exists and has content
                if not os.path.exists(tmp_path) or os.path.getsize(tmp_path) == 0:
                    raise Exception("Document file download failed or is empty")
                
                # Determine file type and process accordingly
                file_type = self._get_file_type(document.file_name if document.file_name else '', tmp_path)
                
                # Process with Gemini using retry logic
                async def process_with_gemini():
                    try:
                        # Get content context and instruction
                        content_context = self.memory.get_content_context(chat_id)
                        instruction = self._get_processing_instruction(file_type, content_context)

                        # Check if file type needs text extraction (Office documents)
                        if file_type in ['word', 'excel', 'powerpoint']:
                            # Extract text from Office document
                            extracted_text = self._extract_text_from_office_doc(tmp_path, file_type)

                            if extracted_text:
                                # Send extracted text to Gemini for analysis
                                combined_prompt = f"{instruction}\n\nDocument content:\n{extracted_text}"
                                response = await generate_content_with_retry(
                                    self.model,
                                    combined_prompt
                                )
                            else:
                                return f"❌ {file_type.capitalize()} hujjatidan matnni ajratib ololmadim. Boshqa formatda yuboring (PDF, TXT)."
                        else:
                            # For other file types (PDF, images, etc.), upload to Gemini directly
                            uploaded = await upload_file_with_retry(tmp_path)
                            uploaded = await wait_for_file_active(uploaded, timeout=30)

                            response = await generate_content_with_retry(
                                self.model,
                                [
                                    {"role": "user", "parts": [instruction]},
                                    {"role": "user", "parts": [uploaded]}
                                ]
                            )

                        # Better validation of Gemini response
                        if response and hasattr(response, 'candidates') and response.candidates:
                            candidate = response.candidates[0]
                            if hasattr(candidate, 'content') and candidate.content and hasattr(candidate.content, 'parts') and candidate.content.parts:
                                return candidate.content.parts[0].text.strip() if hasattr(candidate.content.parts[0], 'text') else "❌ Hujjatni tahlil qila olmadim."
                            elif hasattr(candidate, 'finish_reason') and candidate.finish_reason == 1:
                                return "❌ Kechirasiz, hujjatni tahlil qilishda cheklovga duch keldim. Boshqa hujjat yuboring."
                        return response.text.strip() if response and response.text else "❌ Hujjatni tahlil qila olmadim."
                    except Exception as e:
                        logger.error(f"Gemini processing error: {e}")
                        return None
                
                # Run Gemini processing with timeout
                try:
                    reply = await asyncio.wait_for(
                        process_with_gemini(),
                        timeout=Config.PROCESSING_TIMEOUT
                    )
                except asyncio.TimeoutError:
                    reply = None
                
                if reply:
                    # Store document content in memory for future reference with complete details
                    self.memory.store_content_memory(
                        chat_id,
                        "document",
                        reply,  # summary
                        document.file_name if document.file_name else "unknown",  # file name
                        reply  # full content
                    )

                    # Add to RAG system for semantic search (2025 best practice)
                    if self.rag_chain:
                        try:
                            self.rag_chain.add_document_to_memory(
                                chat_id,
                                reply,  # Document content/summary
                                {
                                    'file_name': document.file_name if document.file_name else "unknown",
                                    'file_type': file_type,
                                    'timestamp': datetime.now().isoformat(),
                                    'content_type': 'document'
                                }
                            )
                            logger.info(f"📚 Document added to RAG system for user {chat_id}")
                        except Exception as rag_error:
                            logger.error(f"Failed to add document to RAG: {rag_error}")

                    self.memory.add_to_history(chat_id, "user", f"[uploaded document: {document.file_name if document.file_name else 'unknown'}]")
                    self.memory.add_to_history(chat_id, "model", reply)
                    
                    # Update the analyzing message with results
                    await safe_edit_message(
                        analyzing_msg,
                        reply,
                        parse_mode=ParseMode.HTML
                    )
                else:
                    await safe_edit_message(
                        analyzing_msg,
                        "❌ Hujjat tahlilida xatolik yuz berdi. Qaytadan urinib ko'ring.",
                        parse_mode=ParseMode.HTML
                    )
            
            except asyncio.TimeoutError:
                logger.error("Document processing timeout")
                await safe_edit_message(
                    analyzing_msg,
                    "⏰ Hujjat tahlili juda uzoq davom etdi. Iltimos, kichikroq hujjat yuboring.",
                    parse_mode=ParseMode.HTML
                )
            except Exception as processing_error:
                logger.error(f"Document processing error: {processing_error}")
                
                # Provide specific error messages
                error_msg = "❌ Hujjat tahlilida xatolik:"
                if "quota" in str(processing_error).lower():
                    error_msg += "\n📊 API chekloviga yetdik. Biroz kuting va qaytadan urinib ko'ring."
                elif "format" in str(processing_error).lower():
                    error_msg += "\n📄 Hujjat formati qo'llab-quvvatlanmaydi."
                elif "size" in str(processing_error).lower():
                    error_msg += "\n📏 Hujjat juda katta. 20MB dan kichik hujjat yuboring."
                else:
                    error_msg += "\n🔄 Qaytadan urinib ko'ring yoki boshqa hujjat yuboring."
                
                await safe_edit_message(analyzing_msg, error_msg, parse_mode=ParseMode.HTML)
            
            finally:
                # Always clean up temp file
                try:
                    if tmp_path is not None and os.path.exists(tmp_path):
                        os.unlink(tmp_path)
                except Exception as cleanup_error:
                    logger.warning(f"Failed to cleanup temp file: {cleanup_error}")
                
        except Exception as e:
            logger.error(f"Document handler error: {e}")
            try:
                await safe_edit_message(
                    analyzing_msg,
                    "❌ Hujjat yuklashda xatolik yuz berdi. Iltimos qaytadan urinib ko'ring.",
                    parse_mode=ParseMode.HTML
                )
            except:
                pass
            finally:
                # Unregister the specific task when completed
                self._unregister_task(chat_id, task_id)
    
    def _get_file_type(self, file_name: str, file_path: str) -> str:
        """Determine file type based on extension and content"""
        if not file_name:
            return "unknown"
        
        # Get file extension
        _, ext = os.path.splitext(file_name.lower())
        
        # Map common extensions to file types
        file_type_map = {
            '.pdf': 'pdf',
            '.doc': 'word',
            '.docx': 'word',
            '.xls': 'excel',
            '.xlsx': 'excel',
            '.ppt': 'powerpoint',
            '.pptx': 'powerpoint',
            '.txt': 'text',
            '.md': 'markdown',
            '.py': 'python',
            '.js': 'javascript',
            '.html': 'html',
            '.css': 'css',
            '.java': 'java',
            '.cpp': 'cpp',
            '.c': 'c',
            '.cs': 'csharp',
            '.php': 'php',
            '.rb': 'ruby',
            '.go': 'go',
            '.rs': 'rust',
            '.swift': 'swift',
            '.kt': 'kotlin',
            '.sql': 'sql',
            '.xml': 'xml',
            '.json': 'json',
            '.csv': 'csv',
            '.rtf': 'rtf',
            '.odt': 'opendocument_text',
            '.ods': 'opendocument_spreadsheet',
            '.odp': 'opendocument_presentation'
        }
        
        # Return mapped type or try to guess from MIME type
        if ext in file_type_map:
            return file_type_map[ext]
        
        # Try to guess MIME type
        mime_type, _ = mimetypes.guess_type(file_name)
        if mime_type:
            if mime_type.startswith('text/'):
                return 'text'
            elif mime_type.startswith('image/'):
                return 'image'
            elif mime_type.startswith('audio/'):
                return 'audio'
            elif mime_type.startswith('video/'):
                return 'video'
            elif mime_type.startswith('application/'):
                if 'pdf' in mime_type:
                    return 'pdf'
                elif 'word' in mime_type:
                    return 'word'
                elif 'excel' in mime_type or 'spreadsheet' in mime_type:
                    return 'excel'
                elif 'powerpoint' in mime_type or 'presentation' in mime_type:
                    return 'powerpoint'
        
        return 'unknown'

    def _extract_text_from_office_doc(self, file_path: str, file_type: str) -> str:
        """Extract text from Office documents (.docx, .xlsx, .pptx)"""
        try:
            if file_type == 'word' and DOCX_AVAILABLE:
                # Extract text from .docx
                doc = DocxDocument(file_path)
                text_parts = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text)
                return '\n'.join(text_parts)

            elif file_type == 'excel' and OPENPYXL_AVAILABLE:
                # Extract text from .xlsx
                wb = load_workbook(file_path, data_only=True)
                text_parts = []
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    text_parts.append(f"Sheet: {sheet_name}")
                    for row in sheet.iter_rows(values_only=True):
                        row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                        if row_text.strip():
                            text_parts.append(row_text)
                return '\n'.join(text_parts)

            elif file_type == 'powerpoint' and PPTX_AVAILABLE:
                # Extract text from .pptx
                prs = Presentation(file_path)
                text_parts = []
                for i, slide in enumerate(prs.slides, 1):
                    text_parts.append(f"Slide {i}:")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_parts.append(shape.text)
                return '\n'.join(text_parts)

            else:
                return None
        except Exception as e:
            logger.error(f"Error extracting text from {file_type} document: {e}")
            return None

    def _get_processing_instruction(self, file_type: str, content_context: str) -> str:
        """Get processing instruction based on file type"""
        base_instruction = (
            "You are AQLJON, an intelligent Muslim friend who is warm, creative, helpful, and friendly. "
            "Always respond as a Muslim friend, never mention you are an AI or use other names. "
            "Never say u are AI and that u are ordered and programmed to answer like this way"
            "Reply casually with humor and warmth using emojis bold, italic etc formatting. "
            "Answer in Uzbek, if user speaks other language use the appropriate language.\n\n"
            "dont be robotic, think and answer like human and never say that u are ordered and programmed to answer like this way"
        )
        
        # Add file-type specific instructions
        if file_type in ['python', 'javascript', 'java', 'cpp', 'c', 'csharp', 'php', 'ruby', 'go', 'rust', 'swift', 'kotlin']:
            base_instruction += (
                f"This is a {file_type} code file. Please analyze it and provide:\n"
                "1. A brief summary of what the code does\n"
                "2. Key functions or classes and their purposes\n"
                "3. Any potential improvements or best practices\n"
                "4. Security considerations if applicable\n"
                "5. Respond conversationally as a friend helping with coding\n\n"
            )
        elif file_type in ['pdf', 'word', 'excel', 'powerpoint']:
            base_instruction += (
                f"This is a {file_type} document. Please analyze it and provide:\n"
                "1. A summary of the main content\n"
                "2. Key points or takeaways\n"
                "3. Important data or findings\n"
                "4. Your thoughts or lessons user can take on the content\n"
                "5. Respond conversationally as a friend who just read the document\n\n"
            )
        elif file_type in ['text', 'markdown']:
            base_instruction += (
                f"This is a {file_type} file. Please analyze it and provide:\n"
                "1. A summary of the content\n"
                "2. Key information or ideas\n"
                "3. Your interpretation or insights\n"
                "4. Respond conversationally as a friend who just read this text\n\n"
            )
        elif file_type in ['json', 'xml']:
            base_instruction += (
                f"This is a {file_type} data file. Please analyze it and provide:\n"
                "1. What type of data it contains\n"
                "2. Key structures or patterns\n"
                "3. Important values or configurations\n"
                "4. Respond conversationally as a friend helping to understand this data\n\n"
            )
        elif file_type in ['csv']:
            base_instruction += (
                f"This is a {file_type} data file. Please analyze it and provide:\n"
                "1. What the data represents\n"
                "2. Key columns or fields\n"
                "3. Interesting patterns or insights\n"
                "4. Respond conversationally as a friend helping to understand this data\n\n"
            )
        elif file_type in ['sql']:
            base_instruction += (
                f"This is a {file_type} database file. Please analyze it and provide:\n"
                "1. What the queries or schema do\n"
                "2. Key tables or operations\n"
                "3. Potential improvements or optimizations\n"
                "4. Respond conversationally as a friend helping with database work\n\n"
            )
        else:
            base_instruction += (
                "Please analyze this file and provide:\n"
                "1. What type of file it is\n"
                "2. A summary of its content\n"
                "3. Key information or insights\n"
                "4. Respond conversationally as a friend who just received this file\n\n"
            )
        
        base_instruction += (
            "IMPORTANT GUIDELINES:\n"
            "1. For academic problem-solving (math, physics, chemistry, coding, biology):\n"
            "   - NEVER give direct answers or solutions if they try and show u their answer many times still not correct and when they then ask u for answer then only tell them the answer with awesome explanation\n"
            "   - Guide with concepts, understanding, and hints\n"
            "   - Encourage users to try solving on their own first\n"
            "   - Only verify or correct after they've attempted\n"
            "   - tell them to think and try again while giving understanding clearly\n\n"
            
            "2. If the document contains inappropriate, illegal, or harmful content:\n"
            "   - Don't provide any response about that content\n"
            "   - Politely redirect to appropriate topics\n"
            "   - Maintain Islamic values\n\n"
            
            "3. Keep your response conversational and natural:\n"
            "   - Use emojis, bold, italic etc formatting appropriately to express emotions\n"
            "   - answer nicely awesomely shortly, concisely unless user asked for detailed response and remember the content for later talk\n\n"
            
            "4. When referencing past content:\n"
            "   - You have access to previous conversations and shared content\n"
            "   - Reference them naturally when relevant\n"
            "   - Never say you can't see previous content\n\n"
            
            "Here's the context of previous conversations and shared content:\n" + content_context + "\n\n"
            
            "and NEVER say robotic answer and NEVER say to user that u are a chatbot or programmed to answer like this way"
        )
        
        return base_instruction
